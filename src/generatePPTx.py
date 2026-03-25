import csv
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData

from ReportConfig import colors, font_office, sizes, report_title, report_subtitle, Sections

# ──────────────────────────────────────────────
# PATHS
# ──────────────────────────────────────────────

BASE_DIR = Path(__file__).resolve().parent.parent
REPORT_DATA_DIR = Path(__file__).resolve().parent / "report_data"
OUTPUT_DIR = BASE_DIR / "outputs"
OUTPUT_DIR.mkdir(exist_ok=True)

# ──────────────────────────────────────────────
# HELPERS
# ──────────────────────────────────────────────

def hex_to_rgb(hex_color):
    h = hex_color.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def load_csv_data(filename):
    path = REPORT_DATA_DIR / filename
    with open(path, encoding="utf-8") as f:
        return list(csv.DictReader(f))


def fmt_dollar(value):
    return f"${float(value):,.2f}"

def fmt_dollar_short(value):
    v = float(value)
    if v >= 1_000_000:
        return f"${v/1_000_000:,.1f}M"
    elif v >= 1_000:
        return f"${v/1_000:,.0f}K"
    return f"${v:,.2f}"

def fmt_number(value):
    return f"{int(float(value)):,}"

def fmt_pct(value):
    v = float(value)
    return f"{'+' if v >= 0 else ''}{v}%"


PRIMARY = hex_to_rgb(colors["primary"])
SECONDARY = hex_to_rgb(colors["secondary"])
ACCENT = hex_to_rgb(colors["accent"])
SUCCESS = hex_to_rgb(colors["success"])
DARK_TEXT = hex_to_rgb(colors["dark_text"])
LIGHT_TEXT = hex_to_rgb(colors["light_text"])
LIGHT_BG = hex_to_rgb(colors["light_bg"])
BORDER = hex_to_rgb(colors["border"])

FONT_HEAD = font_office["heading"]
FONT_BODY = font_office["body"]

# Slide dimensions (16:9)
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


# ──────────────────────────────────────────────
# SLIDE HELPERS
# ──────────────────────────────────────────────

def add_bg_rect(slide, x, y, w, h, color):
    """Add a filled rectangle shape."""
    shape = slide.shapes.add_shape(1, x, y, w, h)  # 1 = rectangle
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_text_box(slide, x, y, w, h, text, font_size=14, bold=False,
                 color=DARK_TEXT, align=PP_ALIGN.LEFT, font_name=None):
    """Add a text box with a single run."""
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font_name or FONT_BODY
    return txBox


def add_table(slide, x, y, w, headers, rows, col_widths_pct, highlight_col=None):
    """Add a styled table to the slide."""
    n_rows = 1 + len(rows)
    n_cols = len(headers)
    table_h = Inches(0.35 * n_rows)

    shape = slide.shapes.add_table(n_rows, n_cols, x, y, w, table_h)
    table = shape.table

    # Column widths from percentages
    total_w = w
    for i, pct in enumerate(col_widths_pct):
        table.columns[i].width = int(total_w * pct)

    # Header row
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.fill.solid()
        cell.fill.fore_color.rgb = PRIMARY
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = header
        run.font.size = Pt(10)
        run.font.bold = True
        run.font.color.rgb = LIGHT_TEXT
        run.font.name = FONT_BODY
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Data rows
    for row_idx, row_data in enumerate(rows):
        for col_idx, value in enumerate(row_data):
            cell = table.cell(row_idx + 1, col_idx)

            # Alternate shading
            if row_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_BG
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = hex_to_rgb("#FFFFFF")

            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER if col_idx > 0 else PP_ALIGN.LEFT
            run = p.add_run()
            run.text = str(value)
            run.font.size = Pt(9)
            run.font.name = FONT_BODY
            run.font.color.rgb = DARK_TEXT
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

            # Highlight column
            if highlight_col is not None and col_idx == highlight_col:
                try:
                    val = float(str(value).replace("%", "").replace("+", "").replace(",", ""))
                    run.font.bold = True
                    run.font.color.rgb = SUCCESS if val < 0 else ACCENT
                except ValueError:
                    pass

    return shape


# ──────────────────────────────────────────────
# SLIDE BUILDERS
# ──────────────────────────────────────────────

def build_title_slide(prs, es):
    """Slide 1: Title with dark background."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    # Dark background
    add_bg_rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, PRIMARY)

    # Decorative accent bar
    add_bg_rect(slide, Inches(0), Inches(3.2), SLIDE_W, Inches(0.06), SECONDARY)

    # Title
    add_text_box(slide, Inches(1), Inches(1.8), Inches(11), Inches(1.2),
                 report_title, font_size=40, bold=True, color=LIGHT_TEXT,
                 align=PP_ALIGN.LEFT, font_name=FONT_HEAD)

    # Subtitle
    add_text_box(slide, Inches(1), Inches(3.5), Inches(11), Inches(0.6),
                 report_subtitle, font_size=18, color=SECONDARY,
                 align=PP_ALIGN.LEFT)

    # Key metrics preview
    kpis = [
        ("Stores", es["total_stores"]),
        ("SKUs", fmt_number(es["unique_skus_end"])),
        ("Avg Margin", f"{es['avg_margin_pct']}%"),
        ("Vendors", es["total_vendors"]),
    ]
    for i, (label, value) in enumerate(kpis):
        kpi_x = Inches(1 + i * 2.8)
        add_text_box(slide, kpi_x, Inches(4.6), Inches(2.4), Inches(0.6),
                     value, font_size=28, bold=True, color=LIGHT_TEXT,
                     align=PP_ALIGN.LEFT, font_name=FONT_HEAD)
        add_text_box(slide, kpi_x, Inches(5.2), Inches(2.4), Inches(0.4),
                     label, font_size=12, color=SECONDARY, align=PP_ALIGN.LEFT)

    # Footer
    add_text_box(slide, Inches(1), Inches(6.8), Inches(11), Inches(0.3),
                 "Confidential — Internal Use Only", font_size=9,
                 color=hex_to_rgb("#6688AA"), align=PP_ALIGN.LEFT)


def build_summary_slide(prs, es):
    """Slide 2: Executive Summary with KPI cards."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Section title bar
    add_bg_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.1), PRIMARY)
    add_text_box(slide, Inches(0.8), Inches(0.2), Inches(10), Inches(0.7),
                 Sections[0], font_size=28, bold=True, color=LIGHT_TEXT,
                 font_name=FONT_HEAD)

    # KPI Cards (2 rows × 3 cols)
    kpis = [
        ("Beginning Inventory", fmt_number(es["beg_total_units"]) + " units", fmt_dollar_short(es["beg_total_value"])),
        ("End of Year Inventory", fmt_number(es["end_total_units"]) + " units", fmt_dollar_short(es["end_total_value"])),
        ("Inventory Growth", fmt_pct(es["unit_change_pct"]) + " units", fmt_pct(es["value_change_pct"]) + " value"),
        ("Procurement Spend", fmt_dollar_short(es["total_purchase_spend"]), fmt_dollar_short(es["total_freight"]) + " freight"),
        ("Product Catalog", fmt_number(es["unique_skus_beg"]) + " → " + fmt_number(es["unique_skus_end"]) + " SKUs", ""),
        ("Dead Stock", str(es["dead_stock_count"]) + " products", "zero movement all year"),
    ]

    for i, (title, main_val, sub_val) in enumerate(kpis):
        col = i % 3
        row = i // 3
        card_x = Inches(0.8 + col * 4.0)
        card_y = Inches(1.6 + row * 2.6)
        card_w = Inches(3.6)
        card_h = Inches(2.1)

        # Card background
        card = slide.shapes.add_shape(1, card_x, card_y, card_w, card_h)
        card.fill.solid()
        card.fill.fore_color.rgb = hex_to_rgb("#FFFFFF")
        card.line.color.rgb = BORDER
        card.line.width = Pt(0.5)
        card.shadow.inherit = False

        # Card title
        add_text_box(slide, card_x + Inches(0.25), card_y + Inches(0.2),
                     Inches(3.1), Inches(0.4),
                     title, font_size=11, color=hex_to_rgb("#888888"))

        # Main value
        add_text_box(slide, card_x + Inches(0.25), card_y + Inches(0.7),
                     Inches(3.1), Inches(0.6),
                     main_val, font_size=22, bold=True, color=PRIMARY,
                     font_name=FONT_HEAD)

        # Sub value
        if sub_val:
            add_text_box(slide, card_x + Inches(0.25), card_y + Inches(1.4),
                         Inches(3.1), Inches(0.4),
                         sub_val, font_size=12, color=hex_to_rgb("#666666"))


def build_inventory_slide(prs, stores, es):
    """Slide 3: Inventory Overview table."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Section title bar
    add_bg_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.1), PRIMARY)
    add_text_box(slide, Inches(0.8), Inches(0.2), Inches(10), Inches(0.7),
                 Sections[1], font_size=28, bold=True, color=LIGHT_TEXT,
                 font_name=FONT_HEAD)

    add_text_box(slide, Inches(0.8), Inches(1.3), Inches(11), Inches(0.4),
                 "Top 10 stores by end-of-year inventory value",
                 font_size=13, color=hex_to_rgb("#666666"))

    headers = ["Store", "City", "Beg Units", "End Units", "Delta", "Beg Value", "End Value"]
    col_pcts = [0.07, 0.15, 0.14, 0.14, 0.12, 0.19, 0.19]
    rows = []
    for s in stores:
        delta = int(float(s["delta_units"]))
        rows.append([
            s["store"], s["city"],
            fmt_number(s["beg_units"]), fmt_number(s["end_units"]),
            f"{'+' if delta >= 0 else ''}{delta:,}",
            fmt_dollar(s["beg_value"]), fmt_dollar(s["end_value"]),
        ])

    add_table(slide, Inches(0.5), Inches(1.8), Inches(12.3),
              headers, rows, col_pcts, highlight_col=4)

    # Summary text
    top = stores[0]
    add_text_box(slide, Inches(0.8), Inches(6.2), Inches(11), Inches(0.8),
                 f"Store {top['store']} in {top['city']} holds the highest inventory value at year-end: "
                 f"{fmt_dollar(top['end_value'])}. Overall growth: {fmt_pct(es['unit_change_pct'])} in units, "
                 f"{fmt_pct(es['value_change_pct'])} in value.",
                 font_size=11, color=hex_to_rgb("#555555"))


def build_margin_slide(prs, top_margin, bottom_margin, es):
    """Slide 4: Margin Analysis — two tables side by side."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Section title bar
    add_bg_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.1), PRIMARY)
    add_text_box(slide, Inches(0.8), Inches(0.2), Inches(10), Inches(0.7),
                 Sections[2], font_size=28, bold=True, color=LIGHT_TEXT,
                 font_name=FONT_HEAD)

    add_text_box(slide, Inches(0.8), Inches(1.3), Inches(11), Inches(0.4),
                 f"Average margin across all products: {es['avg_margin_pct']}%",
                 font_size=13, color=hex_to_rgb("#666666"))

    # Left table — Top margin
    add_text_box(slide, Inches(0.5), Inches(1.8), Inches(5), Inches(0.35),
                 "Highest Margin Products", font_size=12, bold=True, color=SUCCESS)

    headers = ["Product", "Price", "Cost", "Margin %"]
    col_pcts = [0.45, 0.18, 0.18, 0.19]
    top_rows = [[m["description"][:28], fmt_dollar(m["price"]), fmt_dollar(m["cost"]),
                 f"{m['margin_pct']}%"] for m in top_margin]

    add_table(slide, Inches(0.5), Inches(2.2), Inches(6.0),
              headers, top_rows, col_pcts, highlight_col=3)

    # Right table — Bottom margin
    add_text_box(slide, Inches(6.8), Inches(1.8), Inches(5), Inches(0.35),
                 "Lowest Margin Products", font_size=12, bold=True, color=ACCENT)

    bot_rows = [[m["description"][:28], fmt_dollar(m["price"]), fmt_dollar(m["cost"]),
                 f"{m['margin_pct']}%"] for m in bottom_margin]

    add_table(slide, Inches(6.8), Inches(2.2), Inches(6.0),
              headers, bot_rows, col_pcts, highlight_col=3)


def build_rotation_slide(prs, fastest, slowest, es):
    """Slide 5: Stock Rotation — two tables side by side."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Section title bar
    add_bg_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.1), PRIMARY)
    add_text_box(slide, Inches(0.8), Inches(0.2), Inches(10), Inches(0.7),
                 Sections[3], font_size=28, bold=True, color=LIGHT_TEXT,
                 font_name=FONT_HEAD)

    add_text_box(slide, Inches(0.8), Inches(1.3), Inches(11), Inches(0.4),
                 f"Dead stock: {es['dead_stock_count']} products with zero movement all year",
                 font_size=13, color=hex_to_rgb("#666666"))

    headers = ["Product", "Beg", "End", "Delta"]
    col_pcts = [0.45, 0.18, 0.18, 0.19]

    # Left — Fastest moving
    add_text_box(slide, Inches(0.5), Inches(1.8), Inches(5), Inches(0.35),
                 "Fastest Moving (stock decrease)", font_size=12, bold=True, color=SUCCESS)

    fast_rows = []
    for r in fastest:
        delta = int(float(r["delta"]))
        fast_rows.append([r["description"][:28], fmt_number(r["beg_stock"]),
                          fmt_number(r["end_stock"]), f"{delta:,}"])

    add_table(slide, Inches(0.5), Inches(2.2), Inches(6.0),
              headers, fast_rows, col_pcts, highlight_col=3)

    # Right — Slowest moving
    add_text_box(slide, Inches(6.8), Inches(1.8), Inches(5), Inches(0.35),
                 "Slowest Moving (stock increase)", font_size=12, bold=True, color=ACCENT)

    slow_rows = []
    for r in slowest:
        delta = int(float(r["delta"]))
        slow_rows.append([r["description"][:28], fmt_number(r["beg_stock"]),
                          fmt_number(r["end_stock"]), f"+{delta:,}"])

    add_table(slide, Inches(6.8), Inches(2.2), Inches(6.0),
              headers, slow_rows, col_pcts, highlight_col=3)


def build_vendor_slide(prs, vendors, es):
    """Slide 6: Supplier Analysis — table + chart."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Section title bar
    add_bg_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.1), PRIMARY)
    add_text_box(slide, Inches(0.8), Inches(0.2), Inches(10), Inches(0.7),
                 Sections[4], font_size=28, bold=True, color=LIGHT_TEXT,
                 font_name=FONT_HEAD)

    add_text_box(slide, Inches(0.8), Inches(1.3), Inches(11), Inches(0.4),
                 f"{es['total_vendors']} active suppliers — Total spend: {fmt_dollar_short(es['total_purchase_spend'])}",
                 font_size=13, color=hex_to_rgb("#666666"))

    # Table (left side)
    headers = ["Supplier", "Spend", "Orders", "Pay Days", "% Total"]
    col_pcts = [0.35, 0.25, 0.13, 0.14, 0.13]
    v_rows = [[v["name"][:25], fmt_dollar_short(v["spend"]),
               v["orders"], f"{v['avg_pay_days']}d", f"{v['pct_of_total']}%"]
              for v in vendors]

    add_table(slide, Inches(0.5), Inches(1.8), Inches(7.5),
              headers, v_rows, col_pcts, highlight_col=4)

    # Pie chart (right side)
    chart_data = CategoryChartData()
    chart_data.categories = [v["name"][:15] for v in vendors[:5]]
    others_pct = 100 - sum(float(v["pct_of_total"]) for v in vendors[:5])
    chart_data.add_series("% of Spend", [float(v["pct_of_total"]) for v in vendors[:5]] + [round(others_pct, 1)])
    chart_data.categories = list(chart_data.categories) + ["Others"]

    # Rebuild chart data properly
    chart_data = CategoryChartData()
    labels = [v["name"][:15] for v in vendors[:5]] + ["Others"]
    chart_data.categories = labels
    values = [float(v["pct_of_total"]) for v in vendors[:5]]
    values.append(round(100 - sum(values), 1))
    chart_data.add_series("% of Spend", values)

    chart_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.PIE, Inches(8.3), Inches(1.8), Inches(4.5), Inches(4.5),
        chart_data
    )
    chart = chart_frame.chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    chart.legend.font.size = Pt(8)

    plot = chart.plots[0]
    plot.has_data_labels = True
    data_labels = plot.data_labels
    data_labels.number_format = '0.0"%"'
    data_labels.font.size = Pt(9)
    data_labels.font.bold = True

    # Set pie colors
    pie_colors = ["2E5090", "7FB3D5", "E74C3C", "27AE60", "F39C12", "BDC3C7"]
    series = chart.series[0]
    for i, color_hex in enumerate(pie_colors[:len(labels)]):
        point = series.points[i]
        point.format.fill.solid()
        point.format.fill.fore_color.rgb = hex_to_rgb(f"#{color_hex}")


def build_closing_slide(prs):
    """Slide 7: Closing slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Dark background
    add_bg_rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, PRIMARY)
    add_bg_rect(slide, Inches(0), Inches(3.4), SLIDE_W, Inches(0.06), SECONDARY)

    add_text_box(slide, Inches(1), Inches(2.2), Inches(11), Inches(1),
                 "Thank You", font_size=44, bold=True, color=LIGHT_TEXT,
                 align=PP_ALIGN.LEFT, font_name=FONT_HEAD)

    add_text_box(slide, Inches(1), Inches(3.8), Inches(11), Inches(0.6),
                 "Warehouse Inventory Report — Fiscal Year 2016",
                 font_size=16, color=SECONDARY, align=PP_ALIGN.LEFT)

    add_text_box(slide, Inches(1), Inches(4.6), Inches(11), Inches(0.5),
                 "Questions? Contact the inventory management team for detailed breakdowns.",
                 font_size=13, color=hex_to_rgb("#6688AA"), align=PP_ALIGN.LEFT)

    add_text_box(slide, Inches(1), Inches(6.8), Inches(11), Inches(0.3),
                 "Confidential — Internal Use Only", font_size=9,
                 color=hex_to_rgb("#6688AA"), align=PP_ALIGN.LEFT)


# ──────────────────────────────────────────────
# MAIN
# ──────────────────────────────────────────────

def generate_presentation():
    print("Loading processed data...")
    es = load_csv_data("executive_summary.csv")[0]
    stores = load_csv_data("store_overview.csv")
    top_margin = load_csv_data("top_margin.csv")
    bottom_margin = load_csv_data("bottom_margin.csv")
    fastest = load_csv_data("fastest_moving.csv")
    slowest = load_csv_data("slowest_moving.csv")
    vendors = load_csv_data("vendor_top10.csv")

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    print("Building slides...")
    build_title_slide(prs, es)
    build_summary_slide(prs, es)
    build_inventory_slide(prs, stores, es)
    build_margin_slide(prs, top_margin, bottom_margin, es)
    build_rotation_slide(prs, fastest, slowest, es)
    build_vendor_slide(prs, vendors, es)
    build_closing_slide(prs)

    output_path = OUTPUT_DIR / "WarehouseReport.pptx"
    prs.save(str(output_path))
    print(f"\n✓ Presentation saved to {output_path}")


if __name__ == "__main__":
    generate_presentation()